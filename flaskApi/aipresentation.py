#!/usr/bin/env python3
from gptutils import aicomplete, aicomplete_davinci
from pptx import Presentation
from pptx.util import Pt, Inches
import re
import random
import os

system_prompt = """Write a presentation/powerpoint about the user's topic (educational). You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has learning aims (what students will learn from this presentation)
-The presentation has success criteria (what questions students will be able to answer after the presentation)
-The presentation has a summary.
-At least 8 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: LEARNING AIMS
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: SUCCESS CRITERIA
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
...

#Slide: n
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""

if not os.path.exists('GeneratedPresentations'):
        os.makedirs('GeneratedPresentations')
if not os.path.exists('Cache'):
        os.makedirs('Cache')

def create_ppt_text(description, grade ,subject ,number_of_slides):
    user_prompt = f"""The user wants a presentation about "" {description} "" ,grade: {grade}, subject: {subject}, number of slideds: {number_of_slides} """
    response = aicomplete(user_prompt, system_prompt).replace('<br>', '\n')
    return response

def create_ppt(text_file, ppt_name, userid):
    prs = Presentation(f"Designs/Design-3.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for _, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    prs.save(f'GeneratedPresentations/{userid}_{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{userid}_{ppt_name}.pptx"
    return f"{file_path}"

def get_presentation(description ,grade ,subject ,number_of_slides, userid):

    description_string = description
    description_string = re.sub(r'[^\w\s.\-\(\)]', '', description_string)
    description_string = description_string.replace("\n", "")
    with open(f'Cache/{description_string}.txt', 'w', encoding='utf-8') as f:
            f.write(create_ppt_text(description_string, grade ,subject ,number_of_slides))
    file_path = create_ppt(f'Cache/{description_string}.txt', description_string, userid)
    return str(file_path)
